#!/usr/bin/env python3
"""
Presentation Generator POC
Generates branded PowerPoint presentations from a topic using GenAI
"""

import argparse
import os
import sys
from typing import List, Dict
from pptx import Presentation
from pptx.util import Pt
import google.generativeai as genai


class PresentationGenerator:
    def __init__(self, template_path: str, api_key: str = None):
        """Initialize the presentation generator"""
        self.template_path = template_path
        self.prs = None
        
        # Load slide layout configuration
        layout_file = os.path.join(os.path.dirname(__file__), "input/slide_layouts.json")
        try:
            import json
            with open(layout_file, 'r') as f:
                layout_config = json.load(f)
                self.layouts = {layout['name']: layout for layout in layout_config['layouts']}
                self.layouts_by_id = {layout['id']: layout for layout in layout_config['layouts']}
        except FileNotFoundError:
            print(f"❌ Error: slide_layouts.json file not found at: {layout_file}")
            print("Please ensure slide_layouts.json exists in the same directory as this script.")
            sys.exit(1)
        
        # Initialize Google Gemini client
        self.api_key = api_key or os.getenv("GEMINI_API_KEY")
        if not self.api_key:
            print("⚠️  No API key provided. Using mock mode with sample content.")
            self.client = None
        else:
            genai.configure(api_key=self.api_key)
            self.client = genai.GenerativeModel('gemini-2.0-flash')
    
    def generate_outline(self, topic: str) -> Dict:
        """Use GenAI to generate presentation outline and content"""
        print(f"🤖 Generating presentation outline for topic: '{topic}'...")
        
        # Load prompt template from file
        prompt_file = os.path.join(os.path.dirname(__file__), "input/prompt.md")
        try:
            with open(prompt_file, 'r') as f:
                prompt_template = f.read()
            prompt = prompt_template.replace("{topic}", topic)
        except FileNotFoundError:
            print(f"❌ Error: prompt.md file not found at: {prompt_file}")
            print("Please ensure prompt.md exists in the same directory as this script.")
            sys.exit(1)

        if not self.client:
            print("⚠️  Using mock mode (no API key provided)")
            return self._generate_mock_outline(topic)
        
        try:
            response = self.client.generate_content(prompt)
            
            import json
            # Extract JSON from response
            content = response.text
            
            # Save raw response to output directory
            output_dir = "output"
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            safe_topic = "".join(c if c.isalnum() or c in (' ', '_') else '_' for c in topic)
            safe_topic = safe_topic.replace(' ', '_').lower()
            response_file = os.path.join(output_dir, f"{safe_topic}_response.json")
            
            with open(response_file, 'w') as f:
                f.write(content)
            print(f"💾 Saved AI response to: {response_file}")
            
            # Remove markdown code blocks if present
            if content.startswith("```"):
                content = content.split("```")[1]
                if content.startswith("json"):
                    content = content[4:]
            content = content.strip()
            
            outline = json.loads(content)
            # Check if response is in new layouts format
            if 'layouts' in outline:
                print(f"✅ Generated outline with {len(outline.get('layouts', []))} layouts")
            else:
                print(f"✅ Generated outline with {len(outline.get('slides', []))} slides")
            return outline
            
        except Exception as e:
            print(f"❌ Error generating outline: {e}")
            print("Falling back to mock content...")
            return self._generate_mock_outline(topic)
    
    def _generate_mock_outline(self, topic: str) -> Dict:
        """Generate a mock outline when API is not available"""
        return {
            "title": topic,
            "subtitle": "A Comprehensive Overview",
            "slides": [
                {
                    "type": "agenda",
                    "title": "Today's Topics",
                    "points": [
                        "Introduction and Context",
                        "Key Concepts and Principles",
                        "Practical Applications",
                        "Best Practices",
                        "Future Outlook"
                    ]
                },
                {
                    "type": "content",
                    "title": "Introduction",
                    "content": f"This presentation explores {topic}, covering essential concepts, practical applications, and strategic considerations for successful implementation."
                },
                {
                    "type": "content",
                    "title": "Key Concepts",
                    "content": f"Understanding the fundamental principles of {topic} is crucial for effective implementation. This includes:\n• Core definitions and terminology\n• Historical context and evolution\n• Current state and trends\n• Critical success factors"
                },
                {
                    "type": "key_message",
                    "message": f"Success with {topic} requires strategy, execution, and continuous improvement"
                },
                {
                    "type": "content",
                    "title": "Best Practices",
                    "content": "Key recommendations for implementation:\n• Start with clear objectives and metrics\n• Build cross-functional team alignment\n• Adopt iterative, agile approaches\n• Invest in training and change management\n• Monitor progress and adapt continuously"
                },
                {
                    "type": "content",
                    "title": "Looking Ahead",
                    "content": f"The future of {topic} will be shaped by emerging technologies, evolving best practices, and changing organizational needs. Organizations that adapt quickly will gain competitive advantages."
                }
            ]
        }
    
    def create_presentation(self, outline: Dict, output_path: str):
        """Create the presentation using the template and outline"""
        print(f"📄 Creating presentation using template: {self.template_path}")
        
        # Load the branded template
        self.prs = Presentation(self.template_path)
        
        # Check if outline is in new layouts format or old slides format
        if 'layouts' in outline:
            self._create_presentation_from_layouts(outline['layouts'])
        else:
            # Legacy format support
            self._create_presentation_from_slides(outline)
        
        # Save the presentation
        self.prs.save(output_path)
        print(f"✅ Presentation saved to: {output_path}")
        print(f"📊 Total slides: {len(self.prs.slides)}")
    
    def _create_presentation_from_layouts(self, layouts: List[Dict]):
        """Create presentation from new layouts format with embedded content"""
        for layout_data in layouts:
            layout_id = layout_data.get('id')
            layout_name = layout_data.get('name')
            placeholders = layout_data.get('placeholders', [])
            
            # Get the layout from template
            if layout_id not in self.layouts_by_id:
                print(f"⚠️  Warning: Layout ID {layout_id} ({layout_name}) not found in template")
                continue
            
            layout_config = self.layouts_by_id[layout_id]
            slide_layout = self.prs.slide_layouts[layout_id]
            slide = self.prs.slides.add_slide(slide_layout)
            
            # Set content for each placeholder
            for placeholder_data in placeholders:
                idx = placeholder_data.get('idx')
                placeholder_type = placeholder_data.get('type', '')
                content = placeholder_data.get('content')
                max_chars = placeholder_data.get('max_chars')
                
                # Skip if no content or if it's a picture placeholder
                if content is None or 'PICTURE' in placeholder_type:
                    continue
                
                # Set the text in the placeholder
                self._set_placeholder_text(slide, idx, content, max_chars)
    
    def _create_presentation_from_slides(self, outline: Dict):
        """Create presentation from legacy slides format"""
        # Create cover slide
        self._create_cover_slide(outline.get("title", "Presentation"), 
                                outline.get("subtitle", ""))
        
        # Create content slides
        for slide_data in outline.get("slides", []):
            slide_type = slide_data.get("type", "content")
            self._create_slide_by_type(slide_type, slide_data)
        
        # Create closing slide
        self._create_closing_slide()
    
    def _create_slide_by_type(self, slide_type: str, slide_data: Dict):
        """Dynamically create a slide based on type from JSON configuration (legacy support)"""
        # Map slide types to layout names for legacy format
        type_to_layout = {
            "agenda": "agenda_with_image",
            "executive_summary": "executive_summary_with_one_image",
            "key_message": "key_message_02_white_background",
            "content": "content_02_no_image",
            "onepager": "onepager_1"
        }
        
        layout_name = type_to_layout.get(slide_type)
        
        if not layout_name:
            print(f"⚠️  Warning: No layout mapping found for slide type '{slide_type}', skipping...")
            return
        
        # Handle special cases that need custom logic
        if slide_type == "agenda":
            self._create_agenda_slide(slide_data)
        elif slide_type == "onepager":
            self._create_onepager_slide(slide_data)
        else:
            # Use generic layout creation for most slide types
            self._create_slide_from_layout(layout_name, slide_data)
    
    def _set_placeholder_text(self, slide, placeholder_idx: int, text: str, max_chars: int = None):
        """Set text in a placeholder, truncating if necessary"""
        if max_chars and len(text) > max_chars:
            text = text[:max_chars-3] + "..."
        
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == placeholder_idx:
                shape.text = text
                return True
        return False
    
    def _create_slide_from_layout(self, layout_name: str, data: Dict):
        """Generic method to create a slide from layout configuration"""
        if layout_name not in self.layouts:
            print(f"⚠️  Warning: Layout '{layout_name}' not found in configuration")
            return
        
        layout_config = self.layouts[layout_name]
        layout_id = layout_config['id']
        slide_layout = self.prs.slide_layouts[layout_id]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set content for each placeholder based on data
        for placeholder in layout_config['placeholders']:
            if placeholder['type'] == 'text':
                # Map common field names to data
                text = None
                idx = placeholder['idx']
                name = placeholder['name']
                max_chars = placeholder.get('max_chars')
                
                # Try to find matching data
                if 'title' in name.lower() and 'title' in data:
                    text = data['title']
                elif 'subtitle' in name.lower() and 'subtitle' in data:
                    text = data['subtitle']
                elif 'message' in name.lower() and 'message' in data:
                    text = data['message']
                elif 'text' in name.lower() and 'content' in data:
                    text = data['content']
                elif 'presenter' in name.lower() and 'presenter' in data:
                    text = data['presenter']
                elif 'date' in name.lower() and 'date' in data:
                    text = data['date']
                
                if text:
                    self._set_placeholder_text(slide, idx, text, max_chars)
    
    def _create_cover_slide(self, title: str, subtitle: str):
        """Create the cover slide"""
        layout_config = self.layouts['cover']
        slide_layout = self.prs.slide_layouts[layout_config['id']]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Find and set placeholders dynamically
        for placeholder in layout_config['placeholders']:
            if 'title' in placeholder['name'].lower():
                self._set_placeholder_text(slide, placeholder['idx'], title, placeholder.get('max_chars'))
            elif 'subtitle' in placeholder['name'].lower():
                self._set_placeholder_text(slide, placeholder['idx'], subtitle, placeholder.get('max_chars'))
    
    def _create_agenda_slide(self, slide_data: Dict):
        """Create an agenda slide"""
        layout_config = self.layouts['agenda_with_image']
        slide_layout = self.prs.slide_layouts[layout_config['id']]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide_data.get("title", "Agenda")
        for placeholder in layout_config['placeholders']:
            if placeholder['name'] == 'title':
                self._set_placeholder_text(slide, placeholder['idx'], title, placeholder.get('max_chars'))
                break
        
        # Set agenda items
        points = slide_data.get("points", [])
        agenda_placeholders = [p for p in layout_config['placeholders'] if 'agenda_item' in p['name']]
        for i, point in enumerate(points):
            if i < len(agenda_placeholders):
                placeholder = agenda_placeholders[i]
                text = f"{i + 1}. {point}"
                self._set_placeholder_text(slide, placeholder['idx'], text, placeholder.get('max_chars'))
    
    def _create_onepager_slide(self, slide_data: Dict):
        """Create a one-pager summary slide"""
        layout_config = self.layouts['onepager_1']
        slide_layout = self.prs.slide_layouts[layout_config['id']]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide_data.get("title", "Overview")
        for placeholder in layout_config['placeholders']:
            if placeholder['name'] == 'title':
                self._set_placeholder_text(slide, placeholder['idx'], title, placeholder.get('max_chars'))
                break
        
        # Set content in first text placeholder
        content = slide_data.get("content", "")
        text_placeholders = [p for p in layout_config['placeholders'] if p['type'] == 'text' and 'text_' in p['name']]
        if text_placeholders:
            self._set_placeholder_text(slide, text_placeholders[0]['idx'], content, text_placeholders[0].get('max_chars'))
    
    def _create_closing_slide(self):
        """Create a closing/thank you slide"""
        layout_config = self.layouts['salutation']
        slide_layout = self.prs.slide_layouts[layout_config['id']]
        slide = self.prs.slides.add_slide(slide_layout)
        
        for placeholder in layout_config['placeholders']:
            if placeholder['name'] == 'title':
                self._set_placeholder_text(slide, placeholder['idx'], "Thank You", placeholder.get('max_chars'))
                break


def main():
    parser = argparse.ArgumentParser(
        description="Generate a branded PowerPoint presentation from a topic using GenAI"
    )
    parser.add_argument(
        "topic",
        help="The topic for the presentation"
    )
    parser.add_argument(
        "-t", "--template",
        default="input/branding.pptx",
        help="Path to the branded template (default: input/branding.pptx)"
    )
    parser.add_argument(
        "-o", "--output",
        help="Output filename (default: generated_<topic>.pptx)"
    )
    parser.add_argument(
        "-k", "--api-key",
        help="Google Gemini API key (or set GEMINI_API_KEY environment variable)"
    )
    parser.add_argument(
        "-j", "--json",
        help="Use existing JSON response file instead of calling GenAI"
    )
    
    args = parser.parse_args()
    
    # Generate output filename if not provided
    if not args.output:
        safe_topic = "".join(c if c.isalnum() or c in (' ', '_') else '_' for c in args.topic)
        safe_topic = safe_topic.replace(' ', '_').lower()
        args.output = f"output/{safe_topic}.pptx"
    
    # Create output directory if it doesn't exist
    output_dir = os.path.dirname(args.output)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Check if template exists
    if not os.path.exists(args.template):
        print(f"❌ Error: Template file '{args.template}' not found")
        sys.exit(1)
    
    try:
        # Initialize generator
        generator = PresentationGenerator(args.template, args.api_key)
        
        # Generate or load outline
        if args.json:
            # Load existing JSON response
            print(f"📂 Loading existing JSON response from: {args.json}")
            import json
            with open(args.json, 'r') as f:
                content = f.read()
            
            # Remove markdown code blocks if present
            if content.startswith("```"):
                content = content.split("```")[1]
                if content.startswith("json"):
                    content = content[4:]
            content = content.strip()
            
            outline = json.loads(content)
            print(f"✅ Loaded outline with {len(outline.get('slides', []))} slides")
        else:
            # Generate outline using GenAI
            outline = generator.generate_outline(args.topic)
        
        # Create presentation
        generator.create_presentation(outline, args.output)
        
        print(f"\n✨ Success! Your presentation is ready: {args.output}")
        
    except ValueError as e:
        print(f"❌ Configuration Error: {e}")
        print("\nTo use GenAI features, you can:")
        print("  1. Use Google Gemini: Set GEMINI_API_KEY='your-key'")
        print("     Get key at: https://makersuite.google.com/app/apikey")
        print("  2. Or run without API key for mock content mode")
        sys.exit(1)
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()

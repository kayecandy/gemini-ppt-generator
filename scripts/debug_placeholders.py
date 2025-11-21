from pptx import Presentation

prs = Presentation('input/branding.pptx')

print("Analyzing all layouts and their placeholders:\n")
print("=" * 80)

for idx, layout in enumerate(prs.slide_layouts):
    print(f"\nLayout {idx}: {layout.name}")
    print("-" * 80)
    for shape in layout.placeholders:
        print(f"  idx={shape.placeholder_format.idx} | type={shape.placeholder_format.type} | name='{shape.name}'")

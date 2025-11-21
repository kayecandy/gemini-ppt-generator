"""
Analyze the branding.pptx template to understand its structure
"""
from pptx import Presentation

def analyze_template(template_path):
    """Analyze the PowerPoint template and print its structure"""
    prs = Presentation(template_path)
    
    print(f"Template Analysis: {template_path}")
    print("=" * 60)
    
    # Analyze slide dimensions
    print(f"\nSlide Dimensions: {prs.slide_width} x {prs.slide_height}")
    
    # Analyze slide layouts
    print(f"\nAvailable Slide Layouts ({len(prs.slide_layouts)}):")
    for idx, layout in enumerate(prs.slide_layouts):
        print(f"  [{idx}] {layout.name}")
        print(f"      Placeholders: {len(layout.placeholders)}")
        for placeholder in layout.placeholders:
            print(f"        - {placeholder.placeholder_format.type} ({placeholder.name})")
    
    # Analyze existing slides
    print(f"\nExisting Slides: {len(prs.slides)}")
    for idx, slide in enumerate(prs.slides):
        print(f"  Slide {idx + 1}: Layout '{slide.slide_layout.name}'")
    
    # Analyze theme colors
    print("\nTheme Colors:")
    try:
        theme = prs.slide_master.theme
        for idx, color in enumerate(theme.theme_color_scheme):
            print(f"  Color {idx}: {color}")
    except Exception as e:
        print(f"  Could not extract theme colors: {e}")
    
    return prs

if __name__ == "__main__":
    analyze_template("branding.pptx")

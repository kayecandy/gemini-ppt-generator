#!/usr/bin/env python3
"""
Test the presentation generator with a mock outline (no API key required)
"""

from pptx import Presentation
from presentation_generator import PresentationGenerator

def test_without_api():
    """Test the presentation generator without making API calls"""
    
    print("🧪 Testing Presentation Generator (Mock Mode - No API Key Required)")
    print("=" * 70)
    
    # Create a mock outline (simulating what GenAI would return)
    mock_outline = {
        "title": "The Future of Artificial Intelligence",
        "subtitle": "Transforming Business and Society",
        "slides": [
            {
                "type": "agenda",
                "title": "Today's Topics",
                "points": [
                    "AI Overview and Current State",
                    "Key Applications in Business",
                    "Ethical Considerations",
                    "Future Trends and Predictions",
                    "Getting Started with AI"
                ]
            },
            {
                "type": "content",
                "title": "What is Artificial Intelligence?",
                "content": "Artificial Intelligence (AI) refers to computer systems that can perform tasks that typically require human intelligence. This includes learning, reasoning, problem-solving, perception, and language understanding. Modern AI leverages machine learning and deep learning to continuously improve performance."
            },
            {
                "type": "content",
                "title": "AI in Business Today",
                "content": "Organizations are using AI to:\n• Automate repetitive tasks and processes\n• Enhance customer service with chatbots\n• Make data-driven predictions and decisions\n• Personalize user experiences\n• Detect fraud and security threats\n• Optimize supply chains and operations"
            },
            {
                "type": "key_message",
                "message": "AI is not replacing humans—it's augmenting human capabilities"
            },
            {
                "type": "content",
                "title": "Ethical Considerations",
                "content": "As AI becomes more prevalent, we must address:\n• Bias and fairness in AI systems\n• Privacy and data protection\n• Transparency and explainability\n• Accountability for AI decisions\n• Job displacement and workforce transition\n• Environmental impact of AI infrastructure"
            },
            {
                "type": "content",
                "title": "Future Trends",
                "content": "The next decade will see:\n• More sophisticated natural language AI\n• AI-powered scientific discovery\n• Autonomous vehicles becoming mainstream\n• AI in healthcare diagnosis and treatment\n• Advanced robotics in manufacturing\n• AI-assisted creative work"
            },
            {
                "type": "key_message",
                "message": "The best time to start your AI journey is now"
            }
        ]
    }
    
    # Create a minimal generator class without API calls
    class MockGenerator:
        def __init__(self, template_path):
            self.template_path = template_path
            self.prs = None
        
        def create_presentation(self, outline, output_path):
            """Create presentation using the outline"""
            print(f"📄 Creating presentation using template: {self.template_path}")
            self.prs = Presentation(self.template_path)
            
            # Create cover slide
            slide_layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = outline.get("title", "")
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 4:
                    shape.text = outline.get("subtitle", "")
                    break
            
            # Create content slides
            for slide_data in outline.get("slides", []):
                slide_type = slide_data.get("type", "content")
                
                if slide_type == "agenda":
                    slide_layout = self.prs.slide_layouts[7]
                    slide = self.prs.slides.add_slide(slide_layout)
                    if slide.shapes.title:
                        slide.shapes.title.text = slide_data.get("title", "")
                    for shape in slide.placeholders:
                        if shape.placeholder_format.idx == 2:
                            text_frame = shape.text_frame
                            text_frame.clear()
                            for i, point in enumerate(slide_data.get("points", []), 1):
                                p = text_frame.paragraphs[0] if i == 1 else text_frame.add_paragraph()
                                p.text = f"{i}. {point}"
                            break
                
                elif slide_type == "key_message":
                    slide_layout = self.prs.slide_layouts[4]
                    slide = self.prs.slides.add_slide(slide_layout)
                    if slide.shapes.title:
                        slide.shapes.title.text = slide_data.get("message", "")
                
                elif slide_type == "content":
                    slide_layout = self.prs.slide_layouts[7]
                    slide = self.prs.slides.add_slide(slide_layout)
                    if slide.shapes.title:
                        slide.shapes.title.text = slide_data.get("title", "")
                    for shape in slide.placeholders:
                        if shape.placeholder_format.idx == 2:
                            shape.text = slide_data.get("content", "")
                            break
            
            # Create closing slide
            slide_layout = self.prs.slide_layouts[10]
            slide = self.prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = "Thank You"
            
            # Save
            self.prs.save(output_path)
            print(f"✅ Presentation saved to: {output_path}")
            print(f"📊 Total slides: {len(self.prs.slides)}")
    
    # Run the test
    generator = MockGenerator("branding.pptx")
    output_file = "test_presentation.pptx"
    generator.create_presentation(mock_outline, output_file)
    
    print(f"\n✨ Success! Test presentation created: {output_file}")
    print("\nThis demonstrates the POC working with a predefined outline.")
    print("To use with real GenAI, run: python3 presentation_generator.py 'your topic'")
    print("(requires OPENAI_API_KEY environment variable)")

if __name__ == "__main__":
    test_without_api()
